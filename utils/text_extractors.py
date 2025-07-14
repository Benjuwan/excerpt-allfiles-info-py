"""
text_extractors.py
各種テキストファイルからキーワードを抽出する関数群。
ファイル形式ごとに専用の抽出関数を用意し、
extract_text_matchesが拡張子ごとに適切な関数を呼び出す。
"""

from pathlib import Path  # 標準ライブラリ: ファイルパス操作用
import re  # 標準ライブラリ: 正規表現による文字列処理用
import json  # 標準ライブラリ: JSONファイルの読み書き用
import csv  # 標準ライブラリ: CSVファイルの読み書き用
from PyPDF2 import PdfReader  # サードパーティ: PDFファイルからテキスト抽出用
import openpyxl  # サードパーティ: Excel（.xlsx）ファイルの読み書き用
import docx  # サードパーティ: Word（.docx）ファイルの読み書き用
import pptx  # サードパーティ: PowerPoint（.pptx）ファイルの読み書き


def _normalize(text):
    """
    文字列から空白を除去し、小文字に変換する関数
    検索の際に大文字・小文字や空白の違いを無視できるようにします
    """
    return re.sub(r"\s+", "", text).lower()


def extract_txt_md(file_path, norm_keyword):
    """
    .txt, .mdファイルからキーワードを含む行を抽出する関数
    1行ずつ読み込み、キーワードが含まれる行をリストに追加します
    """
    results = []

    # with 文があるので処理終了後は自動的にファイルが閉じる
    # file_path を utf-8 の文字コードで読み込む（開く）際、エンコーディングエラー（デコードできない文字など）が発生した場合は当該ファイルは無視（処理スキップ）する
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        # 読み込んだファイルオブジェクト（f）をイテラブルに、インデックス数値（開始点：1）を用いたループ処理を実施
        for idx, line in enumerate(f, 1):
            # 正規化したキーワードが行に含まれていれば（不要なスペースや改行を取った整形テキストとして）追加
            if norm_keyword in _normalize(line):
                results.append({"line_number": idx, "line": line.strip()})
    return results


def extract_json(file_path, norm_keyword):
    """
    .jsonファイルからキーワードを含む行を抽出する関数
    JSON全体を文字列化し、1行ずつキーワードを検索します
    """
    results = []

    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        try:
            data = json.load(f)
            # JSONを整形して文字列化
            text = json.dumps(data, ensure_ascii=False, indent=2)
            # splitlines() ： 改行コード（\n, \r\n, \r など）を自動的に認識して分割（テキストを1行ずつ処理する）
            for idx, line in enumerate(text.splitlines(), 1):
                if norm_keyword in _normalize(line):
                    results.append({"line_number": idx, "line": line.strip()})
        except Exception as e:
            print(f"[WARN] JSONパース失敗: {file_path}: {e}")
    return results


def extract_csv(file_path, norm_keyword):
    """
    .csvファイルからキーワードを含む行を抽出する関数
    1行ずつ読み込み、カンマ区切りで連結して検索します
    """
    results = []

    # csvファイルの場合、`newline`キーワード引数を指定しないと改行が余分に出力されて各行の間に空行が入ってしまう
    with open(file_path, "r", encoding="utf-8", errors="ignore", newline="") as f:
        reader = csv.reader(f)
        for idx, row in enumerate(reader, 1):
            line = ",".join(row)
            if norm_keyword in _normalize(line):
                results.append({"line_number": idx, "line": line.strip()})
    return results


def extract_pdf(file_path, norm_keyword):
    """
    .pdfファイルからキーワードを含む行を抽出する関数
    各ページごとにテキストを抽出し、行ごとに検索します
    行番号は「pページ番号-l行番号」の形式で返します
    """
    results = []

    try:
        reader = PdfReader(file_path)
        for page_num, page in enumerate(reader.pages, 1):
            text = page.extract_text() or ""

            # 複数ページを持った pdf内の各ページごとのループ処理
            for line_idx, line in enumerate(text.splitlines(), 1):
                if norm_keyword in _normalize(line):
                    results.append(
                        {
                            "line_number": f"p{page_num}-l{line_idx}",
                            "line": line.strip(),
                        }
                    )
    except Exception as e:
        print(f"[WARN] PDF読み込み失敗: {file_path}: {e}")
    return results


def extract_xlsx(file_path, norm_keyword):
    """
    .xlsxファイルからキーワードを含む行を抽出する関数
    全シート・全行を走査し、セルをタブ区切りで連結して検索します
    行番号は「シート名!行番号」の形式で返します
    """
    results = []

    try:
        # file_path のエクセルファイルを読み込み専用として開き、`data_only`キーワード引数によって「セルに数式が入っている場合は数式ではなく計算結果の値を取得」する
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        for sheet in wb.worksheets:
            # シートの各行を値のみで取得してループ処理
            for row_idx, row in enumerate(sheet.iter_rows(values_only=True), 1):
                # 各行のセルの有無を確認し、存在すれば文字列化して返し、それをタブ区切りで連結させる
                line = "\t".join(
                    # リスト内包表記：[式 for 変数 in イテラブル] + 三項演算子：Trueの値 if 条件 else Falseの値
                    [str(cell) if cell is not None else "" for cell in row]
                )
                if norm_keyword in _normalize(line):
                    results.append(
                        {
                            "line_number": f"{sheet.title}!{row_idx}",
                            "line": line.strip(),
                        }
                    )
    except Exception as e:
        print(f"[WARN] Excel読み込み失敗: {file_path}: {e}")
    return results


def extract_docx(file_path, norm_keyword):
    """
    .docxファイルからキーワードを含む行を抽出する関数
    各段落ごとにテキストを抽出し、キーワードを検索します
    """
    results = []

    try:
        doc_file = docx.Document(file_path)
        for idx, para in enumerate(doc_file.paragraphs, 1):
            line = para.text
            if norm_keyword in _normalize(line):
                results.append({"line_number": idx, "line": line.strip()})
    except Exception as e:
        print(f"[WARN] Word読み込み失敗: {file_path}: {e}")
    return results


def extract_pptx(file_path, norm_keyword):
    """
    .pptxファイルからキーワードを含む行を抽出する関数
    各スライド・テキストボックスごとにテキストを抽出し、行ごとに検索します
    行番号は「slideスライド番号-l行番号」の形式で返します
    """
    results = []

    try:
        prs = pptx.Presentation(file_path)
        for slide_idx, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    for line_idx, line in enumerate(shape.text.splitlines(), 1):
                        if norm_keyword in _normalize(line):
                            results.append(
                                {
                                    "line_number": f"slide{slide_idx}-l{line_idx}",
                                    "line": line.strip(),
                                }
                            )
    except Exception as e:
        print(f"[WARN] PowerPoint読み込み失敗: {file_path}: {e}")
    return results


def extract_text_matches(file_path: str, keyword: str) -> list[dict]:
    """
    ファイル拡張子ごとに適切な抽出関数を呼び出し、キーワードを含む行を返すメイン関数
    戻り値: [{line_number: int/str, line: str} ...]
    """
    # suffixプロパティ： ファイルパスとして扱われた要素の拡張子を取得
    ext = Path(file_path).suffix.lower()
    norm_keyword = _normalize(keyword)

    # 拡張子ごとに処理を分岐
    if ext in {".txt", ".md"}:
        return extract_txt_md(file_path, norm_keyword)

    elif ext == ".json":
        return extract_json(file_path, norm_keyword)

    elif ext == ".csv":
        return extract_csv(file_path, norm_keyword)

    elif ext == ".pdf":
        return extract_pdf(file_path, norm_keyword)

    elif ext == ".xlsx":
        return extract_xlsx(file_path, norm_keyword)

    elif ext == ".docx":
        return extract_docx(file_path, norm_keyword)

    elif ext == ".pptx":
        return extract_pptx(file_path, norm_keyword)

    else:
        print(f"[WARN] 未対応拡張子: {file_path}")
        return []
